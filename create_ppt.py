from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Color palette
NAVY = RGBColor(0x1a, 0x1a, 0x2e)
BLUE = RGBColor(0x16, 0x21, 0x3e)
ACCENT = RGBColor(0x0f, 0x3c, 0x78)
GREEN = RGBColor(0x27, 0xae, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)
YELLOW = RGBColor(0xF3, 0x9C, 0x12)

blank_layout = prs.slide_layouts[6]

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_bullet_box(slide, title, bullets, l, t, w, h, bg_color=ACCENT):
    box = add_rect(slide, l, t, w, h, bg_color)
    add_text(slide, title, l+0.15, t+0.1, w-0.3, 0.45, 13, bold=True, color=GREEN)
    y = t + 0.55
    for b in bullets:
        add_text(slide, f"▸  {b}", l+0.15, y, w-0.3, 0.35, 10, color=LIGHT_GRAY)
        y += 0.38

# ─────────────────────────────────────────────
# SLIDE 1 — TITLE
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, NAVY)

# Top accent bar
add_rect(slide, 0, 0, 13.33, 0.08, GREEN)

# Big whale emoji area
add_rect(slide, 0, 0.08, 13.33, 7.42, NAVY)

# Center content
add_text(slide, "🐳", 5.8, 1.2, 1.8, 1.2, 60, align=PP_ALIGN.CENTER)
add_text(slide, "AI Docker NL Health Dashboard", 1.0, 2.5, 11.33, 1.0,
         36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(slide, 4.5, 3.6, 4.3, 0.06, GREEN)
add_text(slide, "Natural Language Powered Container Management", 1.5, 3.8, 10.3, 0.6,
         18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, italic=True)
add_text(slide, "Powered by Groq AI  •  Streamlit  •  Docker SDK  •  Python",
         2.0, 5.2, 9.33, 0.5, 13, color=RGBColor(0x95, 0xa5, 0xa6), align=PP_ALIGN.CENTER)
add_text(slide, "Naveen-300  |  2026", 5.0, 6.6, 3.33, 0.5, 12,
         color=RGBColor(0x7f, 0x8c, 0x8d), align=PP_ALIGN.CENTER)

# Bottom bar
add_rect(slide, 0, 7.3, 13.33, 0.2, GREEN)

# ─────────────────────────────────────────────
# SLIDE 2 — PROBLEM & SOLUTION
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, NAVY)
add_rect(slide, 0, 0, 13.33, 0.08, GREEN)
add_rect(slide, 0, 0.08, 13.33, 1.0, BLUE)
add_text(slide, "Problem & Solution", 0.4, 0.15, 10, 0.8, 28, bold=True, color=WHITE)
add_rect(slide, 0, 7.3, 13.33, 0.2, GREEN)

# Problem box
add_rect(slide, 0.3, 1.3, 5.9, 5.6, ACCENT)
add_text(slide, "❌  The Problem", 0.5, 1.4, 5.5, 0.5, 16, bold=True, color=YELLOW)
problems = [
    "Managing Docker containers requires",
    "memorizing complex CLI commands",
    "No unified health monitoring view",
    "Hard to debug container failures",
    "Non-technical users can't use Docker",
    "Manual log inspection is time-consuming",
    "No natural language interface exists",
]
y = 2.0
for p in problems:
    add_text(slide, f"•  {p}", 0.5, y, 5.5, 0.4, 11, color=LIGHT_GRAY)
    y += 0.42

# Solution box
add_rect(slide, 6.5, 1.3, 6.5, 5.6, RGBColor(0x0d, 0x2b, 0x45))
add_text(slide, "✅  Our Solution", 6.7, 1.4, 6.0, 0.5, 16, bold=True, color=GREEN)
solutions = [
    "Type plain English commands to manage",
    "containers — no CLI knowledge needed",
    "Real-time health dashboard with KPIs",
    "AI-powered ReAct agent (Groq LLM)",
    "7-page interactive Streamlit interface",
    "Full audit trail & prompt logging",
    "GitHub API integration built-in",
]
y = 2.0
for s in solutions:
    add_text(slide, f"•  {s}", 6.7, y, 6.0, 0.4, 11, color=LIGHT_GRAY)
    y += 0.42

# ─────────────────────────────────────────────
# SLIDE 3 — SYSTEM ARCHITECTURE
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, NAVY)
add_rect(slide, 0, 0, 13.33, 0.08, GREEN)
add_rect(slide, 0, 0.08, 13.33, 1.0, BLUE)
add_text(slide, "System Architecture", 0.4, 0.15, 10, 0.8, 28, bold=True, color=WHITE)
add_rect(slide, 0, 7.3, 13.33, 0.2, GREEN)

# Layer boxes
layers = [
    ("🖥️  Frontend Layer", "Streamlit Web UI  •  7 Pages  •  Plotly Charts  •  Custom CSS Theme", GREEN),
    ("🤖  AI Agent Layer", "Groq LLM (llama-3.3-70b)  •  ReAct Loop  •  Natural Language Processing", RGBColor(0x84, 0x29, 0xB0)),
    ("⚙️  Backend Layer", "Docker Manager  •  GitHub API  •  Database (SQLite)  •  LLM Manager", YELLOW),
    ("🐋  Infrastructure", "Docker Engine  •  Docker SDK  •  Container Runtime  •  Docker Socket", RGBColor(0x23, 0x6F, 0xBD)),
]

y = 1.4
for title, desc, color in layers:
    add_rect(slide, 0.4, y, 12.5, 1.1, RGBColor(0x10, 0x28, 0x40))
    add_rect(slide, 0.4, y, 0.12, 1.1, color)
    add_text(slide, title, 0.7, y+0.08, 5.0, 0.45, 13, bold=True, color=color)
    add_text(slide, desc, 0.7, y+0.55, 11.5, 0.45, 11, color=LIGHT_GRAY)
    y += 1.25

# Arrow between layers
add_text(slide, "↕  Data flows between layers via Python modules", 4.0, 6.55, 6.0, 0.4,
         10, color=RGBColor(0x7f, 0x8c, 0x8d), align=PP_ALIGN.CENTER, italic=True)

# ─────────────────────────────────────────────
# SLIDE 4 — KEY FEATURES
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, NAVY)
add_rect(slide, 0, 0, 13.33, 0.08, GREEN)
add_rect(slide, 0, 0.08, 13.33, 1.0, BLUE)
add_text(slide, "Key Features", 0.4, 0.15, 10, 0.8, 28, bold=True, color=WHITE)
add_rect(slide, 0, 7.3, 13.33, 0.2, GREEN)

features = [
    ("🏠 Dashboard", [
        "Live KPI cards",
        "System info",
        "Quick actions",
        "Auto refresh",
    ]),
    ("📦 Containers", [
        "List & filter",
        "Start/Stop/Restart",
        "Status badges",
        "Table view",
    ]),
    ("🤖 AI Agent", [
        "Natural language",
        "ReAct loop steps",
        "7-step reasoning",
        "Groq LLM powered",
    ]),
    ("📋 Logs", [
        "Prompt logs",
        "Audit trail",
        "Container history",
        "SQLite storage",
    ]),
    ("📊 Analytics", [
        "Bar & pie charts",
        "CPU/Memory graphs",
        "Status timeline",
        "Plotly visuals",
    ]),
    ("🔗 GitHub API", [
        "GitHub status",
        "Public events",
        "Rate limit info",
        "Live refresh",
    ]),
]

positions = [(0.3, 1.3), (2.6, 1.3), (4.9, 1.3), (7.2, 1.3), (9.5, 1.3), (11.8, 1.3)]
# 3 per row, 2 rows
positions = [
    (0.3, 1.3), (4.6, 1.3), (8.9, 1.3),
    (0.3, 4.1), (4.6, 4.1), (8.9, 4.1),
]
for i, (title, bullets) in enumerate(features):
    l, t = positions[i]
    add_bullet_box(slide, title, bullets, l, t, 4.0, 2.5)

# ─────────────────────────────────────────────
# SLIDE 5 — AI AGENT (HOW IT WORKS)
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, NAVY)
add_rect(slide, 0, 0, 13.33, 0.08, GREEN)
add_rect(slide, 0, 0.08, 13.33, 1.0, BLUE)
add_text(slide, "AI Agent — How It Works", 0.4, 0.15, 12, 0.8, 28, bold=True, color=WHITE)
add_rect(slide, 0, 7.3, 13.33, 0.2, GREEN)

# ReAct steps
steps = [
    ("1", "User Input", "Type plain English\ne.g. 'restart nginx'"),
    ("2", "LLM Parse", "Groq llama-3.3-70b\nunderstands intent"),
    ("3", "Action Plan", "ReAct loop creates\nstep-by-step plan"),
    ("4", "Execute", "Docker SDK runs\nthe command"),
    ("5", "Observe", "Result captured\nand validated"),
    ("6", "Respond", "Summary shown\nto user"),
]

x = 0.4
for step in steps:
    num, title, desc = step
    add_rect(slide, x, 1.5, 1.9, 2.8, ACCENT)
    add_rect(slide, x, 1.5, 1.9, 0.6, GREEN)
    add_text(slide, num, x, 1.52, 1.9, 0.55, 20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, x+0.05, 2.2, 1.8, 0.5, 12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, desc, x+0.05, 2.75, 1.8, 0.9, 10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    if x < 10.5:
        add_text(slide, "→", x+1.95, 2.6, 0.4, 0.4, 18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    x += 2.2

# Example commands
add_rect(slide, 0.4, 4.6, 12.5, 2.4, RGBColor(0x0d, 0x2b, 0x45))
add_text(slide, "💡  Example Natural Language Commands", 0.6, 4.7, 10, 0.45, 13, bold=True, color=GREEN)
cmds = [
    '"show running containers"   →   Lists all active containers with status',
    '"restart nginx"              →   Finds nginx container and restarts it',
    '"what crashed in last hour"  →   Shows recently exited containers',
    '"get resource usage for mysql"  →   Returns CPU & memory stats',
]
y = 5.2
for cmd in cmds:
    add_text(slide, cmd, 0.6, y, 12.0, 0.35, 10, color=LIGHT_GRAY)
    y += 0.38

# ─────────────────────────────────────────────
# SLIDE 6 — TECH STACK
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, NAVY)
add_rect(slide, 0, 0, 13.33, 0.08, GREEN)
add_rect(slide, 0, 0.08, 13.33, 1.0, BLUE)
add_text(slide, "Technology Stack", 0.4, 0.15, 10, 0.8, 28, bold=True, color=WHITE)
add_rect(slide, 0, 7.3, 13.33, 0.2, GREEN)

tech_categories = [
    ("🖥️  Frontend", GREEN, [
        "Streamlit 1.35+",
        "Plotly 6.x",
        "Custom CSS",
        "Responsive Layout",
    ]),
    ("🤖  AI / LLM", RGBColor(0x84, 0x29, 0xB0), [
        "Groq API",
        "llama-3.3-70b",
        "ReAct Framework",
        "Prompt Engineering",
    ]),
    ("⚙️  Backend", YELLOW, [
        "Python 3.12/3.13",
        "Docker SDK 7.x",
        "SQLite Database",
        "GitHub REST API",
    ]),
    ("🚀  DevOps", RGBColor(0x23, 0x6F, 0xBD), [
        "Docker Engine",
        "Docker Compose",
        "Dockerfile",
        "Health Checks",
    ]),
]

x = 0.3
for cat, color, items in tech_categories:
    add_rect(slide, x, 1.3, 2.9, 5.6, RGBColor(0x0d, 0x2b, 0x45))
    add_rect(slide, x, 1.3, 2.9, 0.55, color)
    add_text(slide, cat, x+0.1, 1.32, 2.7, 0.5, 13, bold=True, color=WHITE)
    y = 2.0
    for item in items:
        add_rect(slide, x+0.15, y, 2.6, 0.5, ACCENT)
        add_text(slide, item, x+0.25, y+0.05, 2.4, 0.4, 11, color=LIGHT_GRAY)
        y += 0.65
    x += 3.18

# ─────────────────────────────────────────────
# SLIDE 7 — CONCLUSION & FUTURE SCOPE
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, NAVY)
add_rect(slide, 0, 0, 13.33, 0.08, GREEN)
add_rect(slide, 0, 0.08, 13.33, 1.0, BLUE)
add_text(slide, "Conclusion & Future Scope", 0.4, 0.15, 12, 0.8, 28, bold=True, color=WHITE)
add_rect(slide, 0, 7.3, 13.33, 0.2, GREEN)

# Achievements
add_rect(slide, 0.3, 1.3, 5.9, 5.5, ACCENT)
add_text(slide, "🏆  What We Achieved", 0.5, 1.4, 5.5, 0.5, 15, bold=True, color=GREEN)
achievements = [
    "Natural language Docker management",
    "7-page interactive web dashboard",
    "AI ReAct agent with step visibility",
    "Real-time container health monitoring",
    "Full audit trail & logging system",
    "GitHub API integration",
    "Docker Compose deployment ready",
    "Clean light theme UI design",
]
y = 2.0
for a in achievements:
    add_text(slide, f"✓  {a}", 0.5, y, 5.5, 0.38, 11, color=LIGHT_GRAY)
    y += 0.4

# Future scope
add_rect(slide, 6.5, 1.3, 6.5, 5.5, RGBColor(0x0d, 0x2b, 0x45))
add_text(slide, "🔭  Future Scope", 6.7, 1.4, 6.0, 0.5, 15, bold=True, color=YELLOW)
future = [
    "Multi-host Docker Swarm support",
    "Kubernetes cluster integration",
    "Email & Slack alert notifications",
    "AI-powered auto-healing containers",
    "User authentication & role management",
    "Custom dashboard widget builder",
    "Mobile-responsive PWA version",
    "CI/CD pipeline integration",
]
y = 2.0
for f in future:
    add_text(slide, f"→  {f}", 6.7, y, 6.0, 0.38, 11, color=LIGHT_GRAY)
    y += 0.4

# Thank you
add_rect(slide, 3.5, 6.85, 6.3, 0.35, GREEN)
add_text(slide, "Thank You  🐳  github.com/Naveen-300/Docker_NL_Health_Dashboard",
         3.5, 6.87, 6.3, 0.3, 10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Save
output_path = r"c:\Users\navee\OneDrive\文档\N Pro (2)\N Pro\docker-nl-dashboard\Docker_NL_Dashboard_Presentation.pptx"
prs.save(output_path)
print(f"PPT saved: {output_path}")
