from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Colors
NAVY    = RGBColor(0x1a, 0x1a, 0x2e)
BLUE    = RGBColor(0x16, 0x21, 0x3e)
ACCENT  = RGBColor(0x0f, 0x3c, 0x78)
GREEN   = RGBColor(0x27, 0xae, 0x60)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xEC, 0xF0, 0xF1)
YELLOW  = RGBColor(0xF3, 0x9C, 0x12)
PURPLE  = RGBColor(0x84, 0x29, 0xB0)
DKBLUE  = RGBColor(0x23, 0x6F, 0xBD)
DARKBG  = RGBColor(0x0d, 0x2b, 0x45)
MGRAY   = RGBColor(0x7f, 0x8c, 0x8d)

blank = prs.slide_layouts[6]

def bg(slide, color=NAVY):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background(); return s

def txt(slide, text, l, t, w, h, size, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb

def header(slide, title, subtitle=None):
    bg(slide)
    rect(slide, 0, 0, 13.33, 0.08, GREEN)
    rect(slide, 0, 0.08, 13.33, 1.05, BLUE)
    txt(slide, title, 0.4, 0.15, 12.5, 0.7, 28, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.4, 0.82, 12.5, 0.35, 13, color=GREEN, italic=True)
    rect(slide, 0, 7.3, 13.33, 0.2, GREEN)

def card(slide, l, t, w, h, title, title_color, lines):
    rect(slide, l, t, w, h, DARKBG)
    rect(slide, l, t, w, 0.55, title_color)
    txt(slide, title, l+0.12, t+0.06, w-0.24, 0.45, 12, bold=True, color=WHITE)
    y = t + 0.62
    for line in lines:
        txt(slide, f"▸  {line}", l+0.12, y, w-0.24, 0.38, 10, color=LGRAY)
        y += 0.38

def bullet_list(slide, items, l, t, w, size=11, color=LGRAY, gap=0.42):
    y = t
    for item in items:
        txt(slide, item, l, y, w, 0.4, size, color=color)
        y += gap

# ══════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, 0, 0, 13.33, 0.08, GREEN)
rect(s, 0, 7.3, 13.33, 0.2, GREEN)
rect(s, 0, 3.55, 13.33, 0.06, GREEN)
txt(s, "🐳", 5.9, 0.9, 1.5, 1.3, 64, align=PP_ALIGN.CENTER)
txt(s, "AI Docker NL Health Dashboard", 0.5, 2.2, 12.3, 1.0,
    38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Intelligent Container Management Through Natural Language",
    1.5, 3.25, 10.3, 0.55, 17, color=GREEN, align=PP_ALIGN.CENTER, italic=True)
txt(s, "Powered by  Groq AI  •  Streamlit  •  Docker SDK  •  Python 3.13",
    2.0, 3.75, 9.33, 0.5, 13, color=LGRAY, align=PP_ALIGN.CENTER)
txt(s, "A Full-Stack AI Dashboard for Real-Time Docker Monitoring & Control",
    1.5, 4.4, 10.33, 0.5, 12, color=MGRAY, align=PP_ALIGN.CENTER, italic=True)
txt(s, "Developed by Naveen-300  |  GitHub: Naveen-300/Docker_NL_Health_Dashboard  |  2026",
    1.0, 6.7, 11.33, 0.4, 11, color=MGRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════
# SLIDE 2 — TABLE OF CONTENTS
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s, "Table of Contents", "Overview of all 20 presentation slides")

toc = [
    ("01", "Title Slide",              "06", "Dashboard Page Deep Dive"),
    ("02", "Table of Contents",        "07", "Container Management Page"),
    ("03", "Project Overview",         "08", "AI Agent — How It Works"),
    ("04", "Problem Statement",        "09", "Logs & Analytics Pages"),
    ("05", "System Architecture",      "10", "GitHub API & Settings"),
    ("11", "Technology Stack",         "16", "Deployment — Docker Compose"),
    ("12", "Groq AI Integration",      "17", "Security & Best Practices"),
    ("13", "ReAct Agent Deep Dive",    "18", "Testing & Demo Results"),
    ("14", "Database & Logging",       "19", "Future Scope & Roadmap"),
    ("15", "UI Design & Theme",        "20", "Conclusion & Thank You"),
]

y = 1.25
for i, (n1, t1, n2, t2) in enumerate(toc):
    rect(s, 0.3, y, 5.9, 0.44, ACCENT)
    rect(s, 6.5, y, 6.5, 0.44, ACCENT)
    rect(s, 0.3, y, 0.55, 0.44, GREEN)
    rect(s, 6.5, y, 0.55, 0.44, GREEN)
    txt(s, n1, 0.3, y+0.04, 0.55, 0.36, 11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, t1, 0.95, y+0.08, 5.1, 0.32, 11, color=LGRAY)
    txt(s, n2, 6.5, y+0.04, 0.55, 0.36, 11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, t2, 7.15, y+0.08, 5.7, 0.32, 11, color=LGRAY)
    y += 0.52

# ══════════════════════════════════════════════════════
# SLIDE 3 — PROJECT OVERVIEW
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s, "Project Overview", "What is AI Docker NL Health Dashboard?")

txt(s, "Project Summary", 0.4, 1.25, 12.5, 0.45, 16, bold=True, color=GREEN)
rect(s, 0.4, 1.7, 12.5, 1.3, DARKBG)
txt(s, "AI Docker NL Health Dashboard is a full-stack web application that allows users to monitor, manage, "
       "and control Docker containers using plain English natural language commands. Instead of memorizing "
       "complex Docker CLI commands, users simply type what they want — and the AI agent handles the rest.",
    0.6, 1.78, 12.1, 1.1, 11, color=LGRAY)

stats = [
    ("7", "Dashboard\nPages"),
    ("6", "Python\nModules"),
    ("22", "Files in\nProject"),
    ("5823", "Lines of\nCode"),
    ("3", "AI\nModels"),
    ("1", "Click\nDeploy"),
]
x = 0.4
for val, label in stats:
    rect(s, x, 3.2, 1.95, 1.3, ACCENT)
    txt(s, val, x, 3.25, 1.95, 0.7, 28, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    txt(s, label, x, 3.9, 1.95, 0.55, 10, color=LGRAY, align=PP_ALIGN.CENTER)
    x += 2.08

txt(s, "Core Objectives", 0.4, 4.65, 12.5, 0.4, 14, bold=True, color=GREEN)
objectives = [
    "✅  Enable non-technical users to manage Docker via natural language",
    "✅  Provide a real-time health dashboard for container monitoring",
    "✅  Build a transparent AI agent with step-by-step execution visibility",
    "✅  Log all operations for audit compliance and debugging",
]
bullet_list(s, objectives, 0.5, 5.1, 12.2, size=11, gap=0.4)

# ══════════════════════════════════════════════════════
# SLIDE 4 — PROBLEM STATEMENT
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s, "Problem Statement", "Why was this project needed?")

rect(s, 0.3, 1.2, 5.9, 5.8, DARKBG)
rect(s, 0.3, 1.2, 5.9, 0.55, RGBColor(0xc0, 0x39, 0x2b))
txt(s, "❌  Current Challenges", 0.5, 1.26, 5.5, 0.45, 14, bold=True, color=WHITE)
problems = [
    "Complex CLI commands hard to remember",
    "No unified real-time health view",
    "Difficult to debug container crashes",
    "Non-technical teams can't use Docker",
    "Manual log inspection is tedious",
    "No natural language interface exists",
    "Scattered tools for different tasks",
    "No audit trail for container actions",
    "Hard to track resource usage trends",
    "Steep learning curve for beginners",
]
y = 1.85
for p in problems:
    txt(s, f"•  {p}", 0.5, y, 5.5, 0.38, 10.5, color=LGRAY)
    y += 0.39

rect(s, 6.5, 1.2, 6.5, 5.8, DARKBG)
rect(s, 6.5, 1.2, 6.5, 0.55, GREEN)
txt(s, "✅  How We Solve It", 6.7, 1.26, 6.1, 0.45, 14, bold=True, color=WHITE)
solutions = [
    "Natural language commands via Groq AI",
    "7-page real-time Streamlit dashboard",
    "AI ReAct agent with visible reasoning",
    "Plain English — no CLI knowledge needed",
    "Automated prompt & audit logging",
    "Unified interface for all Docker ops",
    "Full container history & timeline",
    "SQLite audit trail for all actions",
    "Plotly charts for resource analytics",
    "Simple UI suitable for all skill levels",
]
y = 1.85
for sol in solutions:
    txt(s, f"•  {sol}", 6.7, y, 6.1, 0.38, 10.5, color=LGRAY)
    y += 0.39

# ══════════════════════════════════════════════════════
# SLIDE 5 — SYSTEM ARCHITECTURE
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s, "System Architecture", "4-Layer Architecture Overview")

layers = [
    ("🖥️  Layer 1: Frontend (Presentation)", GREEN,
     "Streamlit Web UI  •  7 Navigation Pages  •  Plotly Interactive Charts  •  Custom CSS Light Theme  •  Responsive Layout"),
    ("🤖  Layer 2: AI Agent (Intelligence)", PURPLE,
     "Groq LLM API  •  llama-3.3-70b-versatile Model  •  ReAct Framework  •  Prompt Engineering  •  Step Tracking"),
    ("⚙️  Layer 3: Backend (Business Logic)", YELLOW,
     "Docker Manager  •  GitHub API Client  •  LLM Manager  •  Database Handler  •  Audit Logger"),
    ("🐋  Layer 4: Infrastructure (Runtime)", DKBLUE,
     "Docker Engine  •  Docker SDK for Python  •  SQLite Database  •  Docker Socket  •  Container Runtime"),
]

y = 1.25
for title, color, desc in layers:
    rect(s, 0.4, y, 12.5, 1.2, DARKBG)
    rect(s, 0.4, y, 0.15, 1.2, color)
    txt(s, title, 0.65, y+0.1, 8.0, 0.45, 13, bold=True, color=color)
    txt(s, desc, 0.65, y+0.62, 12.0, 0.45, 10.5, color=LGRAY)
    if y < 5.5:
        txt(s, "⬇", 6.4, y+1.2, 0.5, 0.35, 14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    y += 1.45

txt(s, "All layers communicate through Python module imports — no microservices overhead, simple and fast.",
    0.4, 7.0, 12.5, 0.25, 10, color=MGRAY, italic=True)

# ══════════════════════════════════════════════════════
# SLIDE 6 — DASHBOARD PAGE
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s, "Dashboard Page — Deep Dive", "The main health overview screen")

rect(s, 0.3, 1.2, 8.2, 5.8, DARKBG)
txt(s, "📊  What the Dashboard Shows", 0.5, 1.28, 7.8, 0.45, 14, bold=True, color=GREEN)
dash_items = [
    "🔢  Total Containers — count of all containers on the system",
    "▶️   Running — containers currently active and serving",
    "⏹️   Stopped — exited or stopped containers",
    "🔄  Restarting — containers in restart loop (potential issues)",
    "⚠️   Unhealthy — containers with CPU usage above 90%",
    "🖥️   Docker Version — engine version info",
    "💻  Operating System — host OS details",
    "🧠  Total Memory — available RAM in GB",
    "🖼️   Images — total Docker images on system",
    "🔄  Refresh Dashboard — live data reload button",
    "📊  View Analytics — quick link to charts page",
]
y = 1.8
for item in dash_items:
    txt(s, item, 0.5, y, 7.8, 0.37, 10.5, color=LGRAY)
    y += 0.4

rect(s, 8.8, 1.2, 4.2, 2.6, ACCENT)
txt(s, "KPI Cards Layout", 8.95, 1.28, 3.9, 0.4, 12, bold=True, color=GREEN)
kpis = ["Total Containers", "Running", "Stopped", "Restarting", "Unhealthy"]
yy = 1.75
for k in kpis:
    rect(s, 8.95, yy, 3.9, 0.36, NAVY)
    txt(s, k, 9.05, yy+0.05, 3.7, 0.28, 10, color=WHITE)
    yy += 0.42

rect(s, 8.8, 4.1, 4.2, 2.9, ACCENT)
txt(s, "Key Implementation", 8.95, 4.18, 3.9, 0.4, 12, bold=True, color=GREEN)
impl = [
    "Uses @st.cache_resource",
    "DockerManager.get_system_info()",
    "DockerManager.list_containers()",
    "st.metric() for KPI cards",
    "st.rerun() for refresh",
]
yy = 4.65
for item in impl:
    txt(s, f"▸ {item}", 8.95, yy, 3.9, 0.37, 10, color=LGRAY)
    yy += 0.42

# ══════════════════════════════════════════════════════
# SLIDE 7 — CONTAINER MANAGEMENT
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s, "Container Management Page", "Full CRUD control over Docker containers")

cols = [
    ("📋  Container Table", GREEN, [
        "Shows all containers in a DataFrame",
        "Columns: Name, Status, Image,",
        "  CPU %, Memory MB, ID",
        "Filter by: all / running /",
        "  exited / paused / restarting",
        "Color-coded status indicators",
        "Sortable and scrollable table",
    ]),
    ("🎮  Container Actions", YELLOW, [
        "Select container from dropdown",
        "▶️  Start — launch stopped container",
        "⏸️  Stop — gracefully stop container",
        "🔄  Restart — restart with same config",
        "Success/Error feedback messages",
        "Real-time status update after action",
        "Uses Docker SDK under the hood",
    ]),
    ("⚙️  Implementation", DKBLUE, [
        "DockerManager.list_containers()",
        "DockerManager.start_container()",
        "DockerManager.stop_container()",
        "DockerManager.restart_container()",
        "pd.DataFrame for table display",
        "st.selectbox for selection",
        "st.success / st.error feedback",
    ]),
]
x = 0.3
for title, color, items in cols:
    card(s, x, 1.2, 4.2, 5.9, title, color, items)
    x += 4.35

# ══════════════════════════════════════════════════════
# SLIDE 8 — AI AGENT OVERVIEW
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s, "AI Agent — Natural Language Control", "How users interact with Docker using plain English")

steps = [
    ("1", "User\nInput", "Type plain\nEnglish command"),
    ("2", "LLM\nParsing", "Groq AI parses\nthe intent"),
    ("3", "Action\nPlan", "ReAct creates\nexecution plan"),
    ("4", "Docker\nExecute", "SDK runs the\nDocker command"),
    ("5", "Result\nCapture", "Output captured\n& validated"),
    ("6", "Summary\nDisplay", "Steps & result\nshown to user"),
]
x = 0.35
for num, title, desc in steps:
    rect(s, x, 1.3, 2.05, 2.8, DARKBG)
    rect(s, x, 1.3, 2.05, 0.6, GREEN)
    txt(s, num, x, 1.32, 2.05, 0.56, 22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, x+0.05, 2.0, 1.95, 0.5, 12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    txt(s, desc, x+0.05, 2.6, 1.95, 0.8, 10, color=LGRAY, align=PP_ALIGN.CENTER)
    if x < 11.0:
        txt(s, "→", x+2.1, 2.45, 0.3, 0.4, 18, bold=True, color=GREEN)
    x += 2.18

rect(s, 0.3, 4.35, 12.7, 2.75, DARKBG)
txt(s, "💬  Natural Language Command Examples", 0.5, 4.42, 10, 0.42, 13, bold=True, color=GREEN)
examples = [
    ('"show running containers"',     "→  Lists all active containers with CPU, memory & status"),
    ('"restart nginx"',               "→  Finds nginx container by name and restarts it safely"),
    ('"what crashed in the last hour"',"→  Shows recently exited containers with timestamps"),
    ('"get resource usage for mysql"', "→  Returns detailed CPU % and Memory MB for mysql"),
    ('"stop redis"',                  "→  Gracefully stops the redis container"),
]
y = 4.9
for cmd, result in examples:
    txt(s, cmd, 0.5, y, 4.5, 0.36, 10, bold=True, color=YELLOW)
    txt(s, result, 5.0, y, 7.8, 0.36, 10, color=LGRAY)
    y += 0.38

# ══════════════════════════════════════════════════════
# SLIDE 9 — REACT AGENT DEEP DIVE
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s, "ReAct Agent — Deep Dive", "Reasoning + Acting framework for AI decision making")

rect(s, 0.3, 1.2, 6.0, 5.8, DARKBG)
txt(s, "🧠  What is ReAct?", 0.5, 1.28, 5.7, 0.42, 14, bold=True, color=GREEN)
react_info = [
    "ReAct = Reasoning + Acting",
    "Each step: Think → Act → Observe",
    "Agent reasons before taking action",
    "Actions are observable and logged",
    "Iterative loop until task complete",
    "Transparent — user sees all steps",
    "Each step has: name, description,",
    "  status, result, duration_ms",
    "Steps displayed in UI expanders",
    "Failed steps shown with ❌ marker",
    "Completed steps shown with ✅",
    "Duration tracked per step (ms)",
]
y = 1.78
for item in react_info:
    txt(s, f"▸  {item}", 0.5, y, 5.6, 0.37, 10.5, color=LGRAY)
    y += 0.38

rect(s, 6.6, 1.2, 6.4, 5.8, DARKBG)
txt(s, "📋  Agent Step Structure", 6.8, 1.28, 6.0, 0.42, 14, bold=True, color=GREEN)
txt(s, "class AgentStep:", 6.8, 1.8, 6.0, 0.35, 10, bold=True, color=YELLOW)
code_lines = [
    "  step_number: int",
    "  name: str",
    "  description: str",
    "  status: pending/running/",
    "         completed/failed",
    "  result: dict | list | str",
    "  start_time: datetime",
    "  end_time: datetime",
    "  get_duration_ms() → float",
    "",
    "Execution Flow:",
    "  1. Parse user intent",
    "  2. Identify Docker action",
    "  3. Validate container name",
    "  4. Execute Docker command",
    "  5. Capture output",
    "  6. Generate summary",
]
y = 2.2
for line in code_lines:
    txt(s, line, 6.8, y, 6.0, 0.32, 10, color=LGRAY)
    y += 0.32

# ══════════════════════════════════════════════════════
# SLIDE 10 — LOGS & ANALYTICS
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s, "Logs & Analytics Pages", "Audit trail, history, and visual insights")

card(s, 0.3, 1.2, 4.0, 2.7, "📋  Prompt Logs Tab", GREEN, [
    "Every AI command is logged",
    "Stores: timestamp, user prompt",
    "Generated action (JSON)",
    "Execution result & status",
    "Duration in milliseconds",
])
card(s, 4.5, 1.2, 4.0, 2.7, "🔍  Audit Logs Tab", YELLOW, [
    "All container operations logged",
    "Who did what and when",
    "DataFrame table display",
    "Filterable and scrollable",
    "SQLite persistent storage",
])
card(s, 8.7, 1.2, 4.3, 2.7, "📅  Container History", DKBLUE, [
    "Timeline of status changes",
    "Container start/stop events",
    "Last 100 events displayed",
    "Timestamp with container name",
    "Useful for post-mortem analysis",
])

rect(s, 0.3, 4.15, 12.7, 2.95, DARKBG)
txt(s, "📊  Analytics Page — Visual Charts", 0.5, 4.22, 10, 0.42, 14, bold=True, color=GREEN)
charts = [
    ("Bar Chart", "Container count grouped by status (running/stopped/paused)"),
    ("Pie Chart", "Percentage distribution of container statuses"),
    ("CPU Chart", "Horizontal bar chart — CPU % per running container"),
    ("Memory Chart", "Memory usage in MB per running container"),
    ("Timeline", "Scatter plot of container status changes over time"),
]
y = 4.72
for chart, desc in charts:
    txt(s, f"📈  {chart}:", 0.5, y, 3.0, 0.36, 10.5, bold=True, color=YELLOW)
    txt(s, desc, 3.6, y, 9.0, 0.36, 10.5, color=LGRAY)
    y += 0.4

# ══════════════════════════════════════════════════════
# SLIDE 11 — GITHUB API & SETTINGS
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s, "GitHub API & Settings Pages", "External integrations and configuration")

card(s, 0.3, 1.2, 6.0, 3.0, "🔗  GitHub API Page", GREEN, [
    "GitHub system status (live)",
    "Incident/outage indicator",
    "API rate limit monitoring",
    "Remaining / Total / Used counts",
    "Progress bar for rate limit",
    "Recent 15 public GitHub events",
    "Event type, repo, actor, timestamp",
])
card(s, 6.6, 1.2, 6.4, 3.0, "⚙️  Settings Page", YELLOW, [
    "Docker socket configuration",
    "Docker host URL setting",
    "Groq API key (password field)",
    "Model selection input",
    "Test Groq connection button",
    "Database path configuration",
    "Log retention days setting",
])

rect(s, 0.3, 4.4, 12.7, 2.7, DARKBG)
txt(s, "🔧  Configuration Details", 0.5, 4.47, 10, 0.42, 14, bold=True, color=GREEN)
config_items = [
    ("GROQ_API_KEY", "Your Groq API key — free at console.groq.com/keys"),
    ("GROQ_MODEL", "llama-3.3-70b-versatile (default) or gemma2-9b-it"),
    ("DOCKER_HOST", "unix:///var/run/docker.sock (default on Linux/Mac)"),
    ("DATABASE_PATH", "data/dashboard.db — SQLite file location"),
    ("LOG_RETENTION", "Days to keep prompt and audit logs (default 30)"),
]
y = 4.95
for key, val in config_items:
    txt(s, f"{key}:", 0.5, y, 3.8, 0.35, 10, bold=True, color=YELLOW)
    txt(s, val, 4.4, y, 8.4, 0.35, 10, color=LGRAY)
    y += 0.38

# ══════════════════════════════════════════════════════
# SLIDE 12 — TECHNOLOGY STACK
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s, "Technology Stack", "All tools, libraries and frameworks used")

tech = [
    ("🖥️  Frontend", GREEN, [
        "Streamlit 1.35+",
        "Plotly 6.x (charts)",
        "Pandas 2.x (tables)",
        "Custom CSS styling",
        "Plotly Express",
        "Plotly Graph Objects",
    ]),
    ("🤖  AI / LLM", PURPLE, [
        "Groq Python SDK 1.4",
        "llama-3.3-70b-versatile",
        "gemma2-9b-it (fallback)",
        "llama-3.1-8b-instant",
        "ReAct framework",
        "Prompt engineering",
    ]),
    ("⚙️  Backend", YELLOW, [
        "Python 3.12 / 3.13",
        "Docker SDK 7.x",
        "SQLite3 (built-in)",
        "Requests 2.31+",
        "JSON / datetime",
        "OS / logging modules",
    ]),
    ("🚀  DevOps", DKBLUE, [
        "Docker Engine 27+",
        "Docker Compose 3.8",
        "Dockerfile (multi-stage)",
        "Health checks",
        "Volume mounts",
        "Bridge networking",
    ]),
]

x = 0.3
for title, color, items in tech:
    card(s, x, 1.2, 3.0, 5.9, title, color, items)
    x += 3.26

# ══════════════════════════════════════════════════════
# SLIDE 13 — GROQ AI INTEGRATION
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s, "Groq AI Integration", "How the LLM powers natural language understanding")

rect(s, 0.3, 1.2, 5.9, 5.8, DARKBG)
txt(s, "🤖  Why Groq?", 0.5, 1.28, 5.6, 0.42, 14, bold=True, color=GREEN)
groq_why = [
    "Extremely fast inference (< 1 second)",
    "Free tier available for development",
    "Multiple powerful open-source models",
    "Simple Python SDK integration",
    "llama-3.3-70b: best reasoning quality",
    "gemma2-9b: fast & lightweight",
    "llama-3.1-8b: ultra-fast responses",
    "JSON-structured output support",
    "Reliable uptime & API availability",
    "No GPU needed on local machine",
    "Supports function/tool calling",
    "Console at console.groq.com/keys",
]
y = 1.78
for item in groq_why:
    txt(s, f"▸  {item}", 0.5, y, 5.6, 0.37, 10.5, color=LGRAY)
    y += 0.37

rect(s, 6.5, 1.2, 6.5, 5.8, DARKBG)
txt(s, "⚙️  LLM Manager — llm.py", 6.7, 1.28, 6.1, 0.42, 14, bold=True, color=GREEN)
llm_details = [
    "class LLMManager:",
    "  __init__(model, api_key)",
    "  generate(prompt) → str",
    "  test_connection() → dict",
    "",
    "System prompt loaded from:",
    "  prompts/system_prompt.txt",
    "",
    "Prompt includes:",
    "  • Available Docker actions",
    "  • Expected JSON output format",
    "  • Container name extraction",
    "  • Error handling instructions",
    "",
    "Response parsed as JSON:",
    "  { action, container, params }",
]
y = 1.78
for line in llm_details:
    color = YELLOW if line.startswith("class") or line.startswith("  __") or line.startswith("  gen") or line.startswith("  test") else LGRAY
    txt(s, line, 6.7, y, 6.1, 0.32, 10.5, color=color)
    y += 0.32

# ══════════════════════════════════════════════════════
# SLIDE 14 — DATABASE & LOGGING
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s, "Database & Logging System", "Persistent storage and audit trail")

card(s, 0.3, 1.2, 3.9, 5.8, "🗄️  SQLite Database", GREEN, [
    "File: data/dashboard.db",
    "3 main tables:",
    "  • prompt_logs",
    "  • audit_logs",
    "  • container_history",
    "Auto-created on startup",
    "No external DB needed",
    "Lightweight & portable",
    "get_prompt_logs(limit=50)",
    "get_audit_logs(limit=50)",
    "get_container_history()",
    "log_prompt() method",
    "log_audit() method",
])
card(s, 4.45, 1.2, 4.0, 5.8, "📝  Prompt Logs Schema", YELLOW, [
    "id: INTEGER PRIMARY KEY",
    "timestamp: DATETIME",
    "user_prompt: TEXT",
    "generated_action: TEXT",
    "execution_result: TEXT",
    "execution_time_ms: REAL",
    "",
    "Stored every AI command",
    "JSON action serialized",
    "Result stored as string",
    "Used for analytics",
    "Displayed in Logs page",
])
card(s, 8.75, 1.2, 4.25, 5.8, "🔍  Audit Log Schema", DKBLUE, [
    "id: INTEGER PRIMARY KEY",
    "timestamp: DATETIME",
    "action: TEXT",
    "container_name: TEXT",
    "status: TEXT",
    "details: TEXT",
    "",
    "Tracks all operations",
    "Start/stop/restart logged",
    "Success and failure both",
    "Retained for 30 days",
    "Exportable as CSV",
])

# ══════════════════════════════════════════════════════
# SLIDE 15 — UI DESIGN & THEME
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s, "UI Design & Theme", "Clean light theme with professional styling")

rect(s, 0.3, 1.2, 7.8, 2.6, DARKBG)
txt(s, "🎨  Design Philosophy", 0.5, 1.28, 7.5, 0.42, 14, bold=True, color=GREEN)
design = [
    "Light background (#f5f7fa) — easy on the eyes for long sessions",
    "White cards with subtle shadows — modern SaaS dashboard feel",
    "Green accent (#27ae60) — primary action and highlight color",
    "Dark navy text (#2c3e50) — high contrast and readable",
    "Streamlit config.toml forces light mode — no dark mode override",
]
y = 1.78
for d in design:
    txt(s, f"▸  {d}", 0.5, y, 7.5, 0.37, 10.5, color=LGRAY)
    y += 0.4

rect(s, 8.4, 1.2, 4.6, 2.6, DARKBG)
txt(s, "🎨  Color Palette", 8.6, 1.28, 4.2, 0.42, 14, bold=True, color=GREEN)
colors = [
    (GREEN,  "#27ae60  Primary Green"),
    (NAVY,   "#1a1a2e  Dark Navy"),
    (ACCENT, "#0f3c78  Card Accent"),
    (LGRAY,  "#ecf0f1  Light Gray"),
    (YELLOW, "#f39c12  Warning/Alt"),
]
yy = 1.78
for col, label in colors:
    rect(s, 8.6, yy, 0.45, 0.32, col)
    txt(s, label, 9.15, yy+0.02, 3.7, 0.3, 10, color=LGRAY)
    yy += 0.4

rect(s, 0.3, 4.05, 12.7, 3.1, DARKBG)
txt(s, "📐  CSS Customizations Applied", 0.5, 4.12, 10, 0.42, 14, bold=True, color=GREEN)
css_items = [
    ("Metric Cards",    "White bg, rounded corners, hover lift shadow, uppercase labels"),
    ("Buttons",         "Green gradient, black text, hover state with deeper green"),
    ("DataFrames",      "White rows, light gray header, subtle row hover highlight"),
    ("Text Inputs",     "White bg, 2px border, green focus ring on active"),
    ("Tabs",            "Green active tab, white inactive, smooth transitions"),
    ("Expanders",       "White bg, border, hover bg change for better UX"),
]
y = 4.62
for comp, desc in css_items:
    txt(s, f"{comp}:", 0.5, y, 2.5, 0.35, 10.5, bold=True, color=YELLOW)
    txt(s, desc, 3.1, y, 9.7, 0.35, 10.5, color=LGRAY)
    y += 0.38

# ══════════════════════════════════════════════════════
# SLIDE 16 — FILE STRUCTURE
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s, "Project File Structure", "Complete codebase organization")

rect(s, 0.3, 1.2, 5.8, 5.8, DARKBG)
txt(s, "📁  Core Python Modules", 0.5, 1.28, 5.5, 0.42, 14, bold=True, color=GREEN)
files = [
    ("app.py",            "Main Streamlit application — all 7 pages"),
    ("agent.py",          "ReactAgent class — ReAct loop implementation"),
    ("llm.py",            "LLMManager — Groq API wrapper"),
    ("docker_manager.py", "DockerManager — Docker SDK wrapper"),
    ("database.py",       "Database — SQLite CRUD operations"),
    ("github_api.py",     "GitHubAPI — REST API client"),
    ("test_demo.py",      "Demo/test script for validation"),
    ("create_ppt.py",     "PPT generation script"),
]
y = 1.78
for fname, desc in files:
    txt(s, fname, 0.5, y, 2.5, 0.36, 10, bold=True, color=YELLOW)
    txt(s, desc, 3.1, y, 2.8, 0.36, 9.5, color=LGRAY)
    y += 0.42

rect(s, 6.4, 1.2, 6.6, 5.8, DARKBG)
txt(s, "📂  Folders & Config Files", 6.6, 1.28, 6.2, 0.42, 14, bold=True, color=GREEN)
folders = [
    (".streamlit/config.toml",  "Theme config — forces light mode"),
    (".env",                    "Environment variables & API keys"),
    (".env.example",            "Template for environment setup"),
    ("requirements.txt",        "Python package dependencies"),
    ("Dockerfile",              "Container build instructions"),
    ("docker-compose.yml",      "Multi-service orchestration"),
    ("data/dashboard.db",       "SQLite database file"),
    ("logs/prompt_log.txt",     "Text-based prompt log file"),
    ("prompts/system_prompt.txt","AI system prompt template"),
    ("start.bat",               "Windows quick-start script"),
    ("start.sh",                "Linux/Mac quick-start script"),
    ("README.md",               "Full documentation"),
]
y = 1.78
for fname, desc in folders:
    txt(s, fname, 6.6, y, 3.0, 0.33, 9.5, bold=True, color=YELLOW)
    txt(s, desc, 9.7, y, 3.1, 0.33, 9.5, color=LGRAY)
    y += 0.37

# ══════════════════════════════════════════════════════
# SLIDE 17 — DEPLOYMENT
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s, "Deployment Guide", "Two ways to run the application")

rect(s, 0.3, 1.2, 5.9, 5.8, DARKBG)
txt(s, "🐍  Option 1: Local Python", 0.5, 1.28, 5.6, 0.42, 14, bold=True, color=GREEN)
txt(s, "Recommended for development", 0.5, 1.75, 5.6, 0.35, 11, color=MGRAY, italic=True)
local_steps = [
    "Step 1: Install dependencies",
    "  pip install -r requirements.txt",
    "",
    "Step 2: Set environment variables",
    "  Copy .env.example to .env",
    "  Add your GROQ_API_KEY",
    "",
    "Step 3: Run the app",
    "  streamlit run app.py",
    "",
    "Step 4: Open browser",
    "  http://localhost:8501",
    "",
    "Requirements:",
    "  • Python 3.12 or 3.13",
    "  • Docker Desktop running",
    "  • Groq API key (free)",
]
y = 2.15
for step in local_steps:
    color = YELLOW if step.startswith("Step") or step.startswith("  pip") or step.startswith("  stream") or step.startswith("  http") or step.startswith("  Copy") or step.startswith("  Add") else LGRAY
    txt(s, step, 0.5, y, 5.6, 0.32, 10, color=color)
    y += 0.32

rect(s, 6.5, 1.2, 6.5, 5.8, DARKBG)
txt(s, "🐋  Option 2: Docker Compose", 6.7, 1.28, 6.1, 0.42, 14, bold=True, color=DKBLUE)
txt(s, "Recommended for production", 6.7, 1.75, 6.1, 0.35, 11, color=MGRAY, italic=True)
docker_steps = [
    "Step 1: Start all services",
    "  docker-compose up -d",
    "",
    "Step 2: Open browser",
    "  http://localhost:8501",
    "",
    "Services started:",
    "  • docker-nl-dashboard (port 8501)",
    "",
    "Volumes mounted:",
    "  • /var/run/docker.sock (read-only)",
    "  • ./data → /app/data",
    "  • ./logs → /app/logs",
    "",
    "Stop everything:",
    "  docker-compose down",
    "",
    "Full cleanup:",
    "  docker-compose down -v",
]
y = 2.15
for step in docker_steps:
    color = YELLOW if step.startswith("Step") or "docker" in step.lower() and step.startswith("  ") else LGRAY
    txt(s, step, 6.7, y, 6.1, 0.32, 10, color=color)
    y += 0.32

# ══════════════════════════════════════════════════════
# SLIDE 18 — SECURITY & BEST PRACTICES
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s, "Security & Best Practices", "How the project handles sensitive data safely")

card(s, 0.3, 1.2, 3.9, 5.8, "🔒  API Key Security", GREEN, [
    ".env file excluded from git",
    ".gitignore covers .env",
    "Password field in UI",
    "Session-only key storage",
    "Never logged or printed",
    "os.getenv() for access",
    "Example file provided",
    "Token in URL cleared after",
    "Regenerate if exposed",
])
card(s, 4.45, 1.2, 4.0, 5.8, "🐋  Docker Security", YELLOW, [
    "Socket mounted read-only",
    "  :ro flag in compose",
    "No privileged mode",
    "Minimal base image",
    "Non-root user in container",
    "Resource limits optional",
    "Network isolation via",
    "  bridge network",
    "Health checks enabled",
])
card(s, 8.75, 1.2, 4.25, 5.8, "📋  Code Best Practices", DKBLUE, [
    "Modular architecture",
    "Separation of concerns",
    "Error handling everywhere",
    "try/except in all API calls",
    "Graceful failure messages",
    "Cache with @st.cache_resource",
    "Type hints in functions",
    "Consistent naming convention",
    "Docstrings on all classes",
])

# ══════════════════════════════════════════════════════
# SLIDE 19 — FUTURE SCOPE
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
header(s, "Future Scope & Roadmap", "Planned enhancements and upcoming features")

phases = [
    ("🔵  Phase 1 — Short Term (3 months)", DKBLUE, [
        "Email & Slack alert notifications for container failures",
        "User authentication with login/logout system",
        "Export logs as CSV / PDF reports",
        "Dark / Light theme toggle in UI",
        "Container resource usage history graphs",
    ]),
    ("🟡  Phase 2 — Medium Term (6 months)", YELLOW, [
        "Docker Swarm multi-host support",
        "Kubernetes cluster integration (kubectl via AI)",
        "Role-based access control (Admin / Viewer)",
        "AI-powered auto-healing (restart on crash)",
        "Custom alert thresholds per container",
    ]),
    ("🟢  Phase 3 — Long Term (12 months)", GREEN, [
        "Mobile-responsive PWA version",
        "CI/CD pipeline integration (GitHub Actions)",
        "Multi-cloud support (AWS ECS, GCP, Azure)",
        "AI-powered capacity planning recommendations",
        "Enterprise SSO and audit compliance reports",
    ]),
]

y = 1.2
for phase_title, color, items in phases:
    rect(s, 0.3, y, 12.7, 1.85, DARKBG)
    rect(s, 0.3, y, 0.15, 1.85, color)
    txt(s, phase_title, 0.55, y+0.08, 8.0, 0.42, 13, bold=True, color=color)
    for i, item in enumerate(items):
        col = 0.55 if i < 3 else 6.8
        row_y = y + 0.55 + (i % 3) * 0.38
        txt(s, f"▸  {item}", col, row_y, 6.0, 0.36, 10, color=LGRAY)
    y += 2.0

# ══════════════════════════════════════════════════════
# SLIDE 20 — CONCLUSION & THANK YOU
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
rect(s, 0, 0, 13.33, 0.08, GREEN)
rect(s, 0, 7.3, 13.33, 0.2, GREEN)

txt(s, "🏆  Conclusion", 0.4, 0.2, 12.5, 0.6, 26, bold=True, color=WHITE)

rect(s, 0.3, 0.95, 5.9, 3.5, DARKBG)
txt(s, "✅  What We Built", 0.5, 1.02, 5.6, 0.42, 13, bold=True, color=GREEN)
built = [
    "Full-stack AI Docker dashboard",
    "Natural language container control",
    "7-page Streamlit web application",
    "ReAct AI agent with step visibility",
    "Real-time health monitoring KPIs",
    "Plotly analytics & visualizations",
    "SQLite audit trail & logging",
    "Docker Compose deployment",
]
y = 1.5
for b in built:
    txt(s, f"✓  {b}", 0.5, y, 5.6, 0.36, 10.5, color=LGRAY)
    y += 0.37

rect(s, 6.5, 0.95, 6.5, 3.5, DARKBG)
txt(s, "💡  Key Learnings", 6.7, 1.02, 6.1, 0.42, 13, bold=True, color=YELLOW)
learnings = [
    "LLMs can bridge technical gaps",
    "ReAct agents improve transparency",
    "Streamlit enables rapid UI building",
    "Docker SDK is powerful & easy",
    "Groq provides fast free inference",
    "Audit logging is critical for trust",
    "Light theme improves usability",
    "Modular code scales easily",
]
y = 1.5
for l in learnings:
    txt(s, f"→  {l}", 6.7, y, 6.1, 0.36, 10.5, color=LGRAY)
    y += 0.37

rect(s, 0.3, 4.6, 12.7, 2.0, ACCENT)
txt(s, "🐳  Thank You!", 0.5, 4.68, 12.3, 0.6, 26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "AI Docker NL Health Dashboard — Making Docker accessible to everyone through AI",
    0.5, 5.3, 12.3, 0.5, 14, color=GREEN, align=PP_ALIGN.CENTER, italic=True)

rect(s, 0.3, 6.15, 5.9, 0.9, DARKBG)
txt(s, "📁  GitHub Repository", 0.5, 6.2, 5.6, 0.35, 11, bold=True, color=GREEN)
txt(s, "github.com/Naveen-300/Docker_NL_Health_Dashboard", 0.5, 6.55, 5.6, 0.35, 10, color=LGRAY)

rect(s, 6.5, 6.15, 6.5, 0.9, DARKBG)
txt(s, "🌐  Live App", 6.7, 6.2, 6.1, 0.35, 11, bold=True, color=GREEN)
txt(s, "http://localhost:8501  •  Powered by Groq AI + Streamlit", 6.7, 6.55, 6.1, 0.35, 10, color=LGRAY)

# Save
out = r"c:\Users\navee\OneDrive\文档\N Pro (2)\N Pro\docker-nl-dashboard\Docker_NL_Dashboard_20Slides.pptx"
prs.save(out)
print(f"Saved: {out}")
